from pptx import Presentation
import os
from hsmolding.const import MOLDFLOW_RESULT,MOLDFLOW_RESULT_CH,MOLDFLOW_RESULT_GROUP
from hsmolding.services.moldflow_report_service import get_moldflow
from django.conf import settings
from pymongo import MongoClient
import imageio
mg_client = MongoClient(settings.MONGO_ADDR)
mg = mg_client[settings.MONGO_DB_HSMOLDING]

class PPT(object):
    def __init__(self, filepath, test_list, company_id, project_id):
        self.filepath = filepath
        self.pptx = Presentation(self.filepath)
        self.test_list = test_list
        self.company_id = company_id
        self.project_id = project_id

    def change_gif_to_mp4(self, image_path, dir_name, file_name):

        # 读取GIF文件
        images = imageio.mimread(image_path)

        # 指定输出MP4文件
        output_mp4_file = file_name+'.mp4'

        # 使用imageio写入MP4文件
        imageio.mimsave(dir_name +"/"+output_mp4_file, images, fps=30)  # 设置适当的帧速率（fps）

    def main(self):
        index = 0
        result = []
        for slide in self.pptx.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                # shape_type 13:图片,14:placeholder, 17:text_box
                if shape.shape_type == 13:
                    plot = {}
                    image = shape.image
                    image_bytes = image.blob
                    dir_name = settings.FILE_STORAGE_PATH + "gsid_"+str(self.company_id) +"/moldflow_report/project_id_"+str(self.project_id)
                    if not os.path.exists(dir_name):
                        os.mkdir(dir_name)
                    if index < len(self.test_list):
                        image_path = dir_name +"/"+str(self.test_list[index])+".gif"
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        animation_url = "gsid_"+str(self.company_id) +"/moldflow_report/project_id_"+str(self.project_id)+"/"+str(self.test_list[index])+".gif",
                        # 如果是在本地服务器,那么注释这一句
                        # self.change_gif_to_mp4(image_path, dir_name, str(self.test_list[index]))
                        # 先拿到之前txt读到的plot
                        filter = {   
                            "project_id":int(self.project_id),
                            "result":{
                                "$elemMatch":{
                                    "name":MOLDFLOW_RESULT_GROUP.get(int(self.test_list[index])),
                                    "children":{
                                        "$elemMatch":{
                                            "id":int(self.test_list[index])
                                        }
                                    }
                                }
                            }                                
                        }
                        one_flow = mg.mold_flow_report_doc.find_one(
                            filter
                        )
                        if type(animation_url) == tuple:
                            animation_url = "".join(animation_url)
                        if one_flow:
                            mg.mold_flow_report_doc.update_many(
                                filter,                            
                                {
                                    "$set": {
                                    "result.$[elem].children.$[subelem].animation_url": animation_url                                  
                                    }
                                },
                                array_filters = [
                                    { "elem.name": MOLDFLOW_RESULT_GROUP.get(int(self.test_list[index]))} ,
                                    { "subelem.id": int(self.test_list[index]) } 
                                    ] ,
                                upsert=True
                            )
                    index = index + 1


# if __name__ == "__main__":
#     ppt = PPT(filepath="C:\\Users\\HUST\\Desktop\\report20230710.pptx", test_list=[
#         1610,
#         1760,
#         1770,
#         1430,
#         1150,
#         1600,
#         1900,
#         1450,
#         1630,
#         1612,
#         1180,
#         1140,
#         1584,
#         1597,
#         7050,
#         1540,
#         1790,
#         1750,
#         1595,
#         1151,
#         1192,
#         1233,
#         1311,
#         1312,
#         1491,
#         1620,
#         1622,
#         1629,
#         1650,
#         1651,
#         1653,
#         1722,
#         1753,
#         7112,
#     ])
#     ppt.main()
